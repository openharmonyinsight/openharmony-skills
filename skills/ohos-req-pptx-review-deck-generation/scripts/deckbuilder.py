# -*- coding: utf-8 -*-
"""
deckbuilder.py — a batteries-included helper for building clean, consistent
16:9 PowerPoint decks with python-pptx.

LIGHT THEME variant
-------------------
* No top color band — each page is a title + thin accent underline on white.
* Slide titles are 32pt Microsoft YaHei.
* OpenHarmony logo sits in the bottom-left corner of every page.
* System diagrams are real drawing boxes (rounded rectangles) joined by arrows.
* Overall palette is light; heavy fills are avoided.

You only supply CONTENT (titles, bullet lists, table rows, diagram stages).
This module owns all layout, coordinates, colors, fonts, and spacing.

Public API is unchanged from the original deckbuilder:
    Deck().cover / content_slide / bullets_slide / table_slide /
           flow_slide / layered_diagram_slide / architecture_slide / save
"""

import os
import math
import warnings
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --------------------------------------------------------------------------
# palette — pass colors by NAME (string) anywhere a color is accepted
# theme "red ink": near-black ink titles (素); ONE red accent carries the
#   conclusion + structure highlight; soft grey for chrome; blue for value-page
#   body text + table headers; amber reserved for change points.
# --------------------------------------------------------------------------
PALETTE = {
    "primary": RGBColor(0x23, 0x28, 0x2E),   # near-black ink — titles (subdued)
    "accent":  RGBColor(0xD8, 0x1E, 0x1E),   # red — conclusion/underline/headers/primary arrows
    "dark":    RGBColor(0x1A, 0x1D, 0x21),   # body text — soft black
    "ink":     RGBColor(0x23, 0x28, 0x2E),
    "grey":    RGBColor(0x8A, 0x90, 0x98),   # soft neutral grey (subtitles/structure)
    "light":   RGBColor(0xF5, 0xF6, 0xF7),   # card background (neutral off-white)
    "white":   RGBColor(0xFF, 0xFF, 0xFF),
    "green":   RGBColor(0x3A, 0x9A, 0x57),   # logo green
    "orange":  RGBColor(0xBE, 0x7A, 0x22),   # amber → change border/badge/title only
    "red":     RGBColor(0xB5, 0x40, 0x3A),   # muted brick-red → genuine risk only
    "blue":    RGBColor(0x2B, 0x5A, 0xA6),   # calm slate-blue (value body + table-header family)
}
# soft tint backgrounds keyed by accent name (kept very light, desaturated)
_TINT = {
    "primary": RGBColor(0xEC, 0xED, 0xEF),
    "accent":  RGBColor(0xF8, 0xE0, 0xE0),   # light red (table total / accent cards)
    "green":   RGBColor(0xEA, 0xF3, 0xEC),
    "orange":  RGBColor(0xFA, 0xE9, 0xCF),   # light amber (change box fill)
    "red":     RGBColor(0xF4, 0xE6, 0xE5),
    "grey":    RGBColor(0xEF, 0xEE, 0xEC),   # warm neutral (coordinates under the red lead)
}
# neutral box fill / border for un-changed diagram nodes (warm-leaning, airy)
_BOX_FILL = RGBColor(0xF8, 0xF7, 0xF6)
_BOX_LINE = RGBColor(0xCB, 0xC8, 0xC6)
_RULE     = RGBColor(0xE8, 0xE5, 0xE3)        # faint warm separators


def _color(c):
    """Accept an RGBColor or a palette name string."""
    if isinstance(c, RGBColor):
        return c
    return PALETTE.get(str(c).lower(), PALETTE["primary"])


def _tint(name):
    return _TINT.get(str(name).lower(), PALETTE["light"])


def _num(v):
    """Parse the leading number out of a value like '1200 行' / '2.5 人月' /
    '1,200 行' (thousands separator). Returns None when nothing numeric is
    present, so callers can distinguish a genuine 0 from 'unknown' instead of
    silently folding missing data into the total."""
    if v is None:
        return None
    if isinstance(v, (int, float)):
        return float(v)
    import re
    s = str(v).strip()
    if not s or s == Deck._TBD:
        return None
    m = re.search(r"-?[\d,]*\.?\d+", s)
    if not m:
        return None
    digits = m.group().replace(",", "")
    try:
        return float(digits)
    except ValueError:
        return None


# Titles that MUST use the dedicated two-column builders, never the generic
# card/bar grid. A weak model tends to reach for content_slide (four boxes) for
# these — we warn (not fail) so the mistake is visible without breaking other use.
_TWO_COL_TITLES = ("需求价值描述", "需求设计方案")


def _warn_wrong_builder(title, used, suggested):
    t = str(title or "")
    for key in _TWO_COL_TITLES:
        if key in t:
            warnings.warn(
                "deckbuilder: 页面 '%s' 用了 %s()，但需求评审的『%s』页应改用 "
                "%s()（左文右图两栏；设计页右侧传 diagram= 画框图）。"
                "最省心的做法是整套评审直接调 Deck.requirement_review_deck(spec)，"
                "由代码固定 8 页结构与每页 builder。"
                % (t, used, key, suggested),
                stacklevel=3)
            break




class Deck:
    # 16:9 canvas
    W = Inches(13.333)
    H = Inches(7.5)
    HEADER_H = Inches(1.15)          # (retained for compatibility; unused)
    BODY_TOP = Inches(1.62)          # first usable y below the title block
    MARGIN = Inches(0.5)
    BODY_BOTTOM = Inches(6.82)       # last usable y (footer strip below)
    LOGO_RATIO = 4.723               # width/height of oh_logo.png

    def __init__(self, font="Microsoft YaHei", logo="oh_logo.png"):
        self.font = font
        self.prs = Presentation()
        self.prs.slide_width = self.W
        self.prs.slide_height = self.H
        self._blank = self.prs.slide_layouts[6]
        self._page = 0               # running page number (cover excluded)
        # resolve logo relative to cwd or this file's directory
        self._logo = None
        if logo:
            for cand in (logo, os.path.join(os.path.dirname(__file__), logo)):
                if os.path.exists(cand):
                    self._logo = cand
                    break

    # ---------------- low-level primitives ----------------
    def _slide(self):
        return self.prs.slides.add_slide(self._blank)

    def _run(self, r, size, color, bold=False, italic=False):
        r.font.size = Pt(size)
        r.font.color.rgb = _color(color)
        r.font.bold = bold
        r.font.italic = italic
        if self.font:
            r.font.name = self.font

    def _textbox(self, slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        return tf

    def _rect(self, slide, l, t, w, h, fill, line=None, line_w=0.75,
              shape=MSO_SHAPE.RECTANGLE):
        sp = slide.shapes.add_shape(shape, l, t, w, h)
        sp.fill.solid()
        sp.fill.fore_color.rgb = _color(fill)
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = _color(line)
            sp.line.width = Pt(line_w)
        sp.shadow.inherit = False
        return sp

    def _arrow(self, slide, l, t, w, h, shape, color="accent"):
        sp = slide.shapes.add_shape(shape, l, t, w, h)
        sp.fill.solid()
        sp.fill.fore_color.rgb = _color(color)
        sp.line.fill.background()
        sp.shadow.inherit = False
        return sp

    def _footer(self, slide):
        """OpenHarmony logo, bottom-left, on every page."""
        if not self._logo:
            return
        h = Inches(0.30)
        w = Emu(int(int(h) * self.LOGO_RATIO))
        slide.shapes.add_picture(self._logo, self.MARGIN,
                                 self.H - h - Inches(0.16), width=w, height=h)

    def _header(self, slide, title, subtitle=None, takeaway=None, numbered=True):
        """Light title block: no band — title + thin accent underline.
        Page number sits top-right; OpenHarmony logo goes bottom-left.
        `takeaway` is the page's one-line conclusion, rendered bold-accent right
        under the title so every page leads with its point. It is REQUIRED — a
        content page without a conclusion has no 突出重点, so the build aborts."""
        if not (takeaway and str(takeaway).strip()):
            raise ValueError(
                "deckbuilder: slide %r was built without takeaway=. Every content "
                "page MUST lead with its conclusion so the page 突出重点. Pass "
                "takeaway=\"结论：…\" — a one-line VERDICT (an assertion, not a "
                "restatement of the title), e.g. "
                "takeaway=\"结论：不改变公开 API 行为，应用无需适配\". "
                "(subtitle= is only a neutral scope note and does not satisfy this.)"
                % (title or "")
            )
        x = self.MARGIN
        if numbered:
            self._page += 1
            tf = self._textbox(slide, self.W - Inches(1.4), Inches(0.34),
                               Inches(0.9), Inches(0.5), MSO_ANCHOR.MIDDLE)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
            r = p.add_run(); r.text = "%02d" % self._page
            self._run(r, 14, "grey", bold=True)
        # title — single line, auto-shrink if needed
        title_w = self.W - x - Inches(1.5)
        tf = self._textbox(slide, x, Inches(0.36),
                           title_w, Inches(0.74),
                           MSO_ANCHOR.MIDDLE)
        tf.word_wrap = False
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = title
        tsize = self._fit_size(title, 32, title_w / 914400.0, minsize=18)
        self._run(r, tsize, "primary", bold=True)
        # accent underline
        self._rect(slide, x, Inches(1.16), Inches(2.4), Inches(0.055), "accent")
        self._rect(slide, x + Inches(2.4), Inches(1.175), self.W - x - Inches(2.4)
                   - Inches(0.5), Inches(0.025), _RULE)
        if takeaway:
            # small red kicker square + bold red one-line conclusion —
            # the verdict is the page's highlight; ink title stays subdued above it
            self._rect(slide, x, Inches(1.335), Inches(0.09), Inches(0.09), "accent")
            tf = self._textbox(slide, x + Inches(0.17), Inches(1.235),
                               self.W - x - Inches(0.77), Inches(0.34))
            tf.word_wrap = False
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = takeaway
            tw = self.W - x - Inches(0.77)
            tsize = self._fit_size(takeaway, 14, tw / 914400.0, minsize=11)
            self._run(r, tsize, "accent", bold=True)
        elif subtitle:
            tf = self._textbox(slide, x, Inches(1.24),
                               self.W - x - Inches(0.6), Inches(0.32))
            r = tf.paragraphs[0].add_run(); r.text = subtitle
            self._run(r, 12.5, "grey")
        self._footer(slide)

    def _card(self, slide, l, t, w, h, title, bullets, accent="accent"):
        self._rect(slide, l, t, w, h, _tint(accent), line=_RULE, line_w=0.75)
        self._rect(slide, l, t, Inches(0.09), h, accent)      # left accent stripe
        tf = self._textbox(slide, l + Inches(0.22), t + Inches(0.13),
                           w - Inches(0.36), h - Inches(0.26))
        p = tf.paragraphs[0]; p.space_after = Pt(5)
        r = p.add_run(); r.text = title; self._run(r, 15, "dark", bold=True)
        for ln in bullets:
            pp = tf.add_paragraph(); pp.space_after = Pt(3)
            r = pp.add_run(); r.text = "• " + str(ln)
            self._run(r, 12.5, "dark")

    @staticmethod
    def _fit_size(text, base, inner_w_in, minsize=8.5):
        """Shrink a single-line label's font so it fits inner_w_in inches.
        CJK glyphs ~ base pt wide; latin/space ~ 0.55*base."""
        units = sum(1.0 if ord(c) > 0x2E80 else 0.55 for c in text)
        width_in = units * base / 72.0
        if width_in <= inner_w_in or width_in <= 0:
            return base
        return max(minsize, base * inner_w_in / width_in)

    @staticmethod
    def _text_width(text, pt):
        """Estimated width in inches for a single line of text at given pt."""
        units = sum(1.0 if ord(c) > 0x2E80 else 0.55 for c in text)
        return units * pt / 72.0

    @staticmethod
    def _text_units(text):
        """Approx glyph-width units for a line (CJK=1.0, latin/space=0.55)."""
        return sum(1.0 if ord(c) > 0x2E80 else 0.55 for c in str(text))

    @classmethod
    def _wrapped_rows(cls, text, pt, w_in):
        """Estimate how many wrapped rows `text` occupies at `pt` size within
        `w_in` inches of column width, honoring explicit newlines. This is
        what table_slide uses to size rows to REAL cell content instead of a
        flat row/column-count guess."""
        text = str(text)
        if not text:
            return 1
        cap = max(1.0, w_in * 72.0 / pt)
        total = 0
        for line in text.split("\n"):
            if not line:
                total += 1
                continue
            total += max(1, math.ceil(cls._text_units(line) / cap))
        return total

    @classmethod
    def _fit_box_fonts(cls, title, lines, inner_w_in, inner_h_in,
                       title_pt, sub_pt, min_title=8.0, min_sub=7.0):
        """Shrink a box's title+body fonts so the WRAPPED text fits inside
        (inner_w_in × inner_h_in). Estimates wrapped-row count at each size and
        scales down until the stacked height fits or the floor is hit. This is
        what keeps long node text from spilling out of a diagram box."""
        if inner_w_in <= 0 or inner_h_in <= 0:
            return title_pt, sub_pt
        lines = [str(l) for l in (lines or [])]

        def rows_at(t_pt, s_pt):
            # chars-per-row is width-driven; a row holds inner_w worth of units
            def rows(text, pt):
                if not text:
                    return 1
                cap = max(1.0, inner_w_in * 72.0 / pt)      # units per row
                import math
                return max(1, math.ceil(cls._text_units(text) / cap))
            total = rows(title, t_pt)
            for ln in lines:
                total_rows_ln = rows(ln, s_pt)
                total += total_rows_ln
            return total

        def height_at(t_pt, s_pt):
            # first row uses title line-height, remaining rows use body height
            r_title = 1  # title occupies at least one row of title height
            def rows(text, pt):
                if not text:
                    return 1
                cap = max(1.0, inner_w_in * 72.0 / pt)
                import math
                return max(1, math.ceil(cls._text_units(text) / cap))
            h = rows(title, t_pt) * t_pt * 1.32 / 72.0
            for ln in lines:
                h += rows(ln, s_pt) * s_pt * 1.28 / 72.0
            return h

        t_pt, s_pt = float(title_pt), float(sub_pt)
        # up to ~14 shrink steps of 7% each
        for _ in range(14):
            if height_at(t_pt, s_pt) <= inner_h_in:
                break
            if t_pt <= min_title and s_pt <= min_sub:
                break
            t_pt = max(min_title, t_pt * 0.93)
            s_pt = max(min_sub, s_pt * 0.93)
        return t_pt, s_pt

    @classmethod
    def _measure_boxes(cls, nodes, title_pt=12.5, sub_pt=10.5,
                       pad_w=0.28, pad_h=0.24):
        """Compute uniform (w, h) that fits all nodes' text content.
        Returns (box_w, box_h, title_pt, sub_pt) — font sizes may shrink."""
        max_w = 0
        max_lines = 0
        for nd in nodes:
            title = nd.get("title", "")
            lines = nd.get("lines", [])
            tw = cls._text_width(title, title_pt)
            for ln in lines:
                tw = max(tw, cls._text_width(str(ln), sub_pt))
            max_w = max(max_w, tw)
            max_lines = max(max_lines, 1 + len(lines))
        box_w = max_w + pad_w
        lh_title = title_pt * 1.45 / 72.0
        lh_sub = sub_pt * 1.4 / 72.0
        box_h = lh_title + lh_sub * max(0, max_lines - 1) + pad_h
        return (box_w, box_h, title_pt, sub_pt)

    def _flowbox(self, slide, l, t, w, h, title, lines, change=False,
                 title_pt=12.5, sub_pt=10.5):
        """A real drawing box (rounded rectangle) for system diagrams.
        Title AND sub-lines auto-shrink together so wrapped text stays inside
        the box (no more text spilling past the border on verbose content)."""
        border = "orange" if change else _BOX_LINE
        fill = _tint("orange") if change else _BOX_FILL
        sp = self._rect(slide, l, t, w, h, fill, line=border,
                        line_w=2.25 if change else 1.25,
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = self._textbox(slide, l + Inches(0.08), t + Inches(0.08),
                           w - Inches(0.16), h - Inches(0.14), MSO_ANCHOR.MIDDLE)
        inner_w = (w - Inches(0.16)) / 914400.0           # box inner width (in)
        inner_h = (h - Inches(0.14)) / 914400.0           # box inner height (in)
        # shrink title+body together to fit the wrapped text into the box
        tsize, ssize = self._fit_box_fonts(title, lines, inner_w, inner_h,
                                           title_pt, sub_pt)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title
        self._run(r, tsize, _color("orange") if change else "dark",
                  bold=True)
        for sub in lines:
            pp = tf.add_paragraph(); pp.alignment = PP_ALIGN.CENTER
            pp.space_before = Pt(1)
            r = pp.add_run(); r.text = str(sub); self._run(r, ssize, "dark")
        if change:
            bw, bh = Inches(0.52), Inches(0.23)
            bx = l + w - bw - Inches(0.04)
            by = t - Inches(0.12)                 # straddle the top-right corner
            self._rect(slide, bx, by, bw, bh, "orange",
                       shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            tf2 = self._textbox(slide, bx, by, bw, bh, MSO_ANCHOR.MIDDLE)
            bp = tf2.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
            r = bp.add_run(); r.text = "★变更"; self._run(r, 8, "white", bold=True)
        return sp

    @staticmethod
    def _cell_border(cell, color="C2CCD8", w_pt=0.75):
        """Solid line on all four sides of a table cell (grid separators)."""
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        emu = str(int(w_pt * 12700))
        for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
            for el in tcPr.findall(qn(tag)):
                tcPr.remove(el)
        # insert reversed so final child order is lnL, lnR, lnT, lnB (before fill)
        for tag in ("a:lnB", "a:lnT", "a:lnR", "a:lnL"):
            ln = tcPr.makeelement(qn(tag), {"w": emu, "cap": "flat",
                                            "cmpd": "sng", "algn": "ctr"})
            sf = ln.makeelement(qn("a:solidFill"), {})
            sf.append(ln.makeelement(qn("a:srgbClr"), {"val": color}))
            ln.append(sf)
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": "solid"}))
            tcPr.insert(0, ln)

    def _style_table(self, table, header_size=15, body_size=13.5,
                     highlight_last=False):
        HDR = RGBColor(0xC6, 0xD7, 0xEC)      # light slate-blue header (same family as body blue)
        GRID = "000000"                        # black solid grid lines
        for cell in table.rows[0].cells:
            cell.fill.solid(); cell.fill.fore_color.rgb = HDR
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            self._cell_border(cell, color=GRID, w_pt=0.75)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    self._run(r, header_size, "dark", bold=True)
        n = len(table.rows)
        for ri in range(1, n):
            last = highlight_last and ri == n - 1
            for cell in table.rows[ri].cells:
                if last:
                    cell.fill.solid(); cell.fill.fore_color.rgb = _TINT["accent"]
                else:
                    cell.fill.background()
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                self._cell_border(cell, color=GRID, w_pt=0.75)
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        self._run(r, body_size, "dark", bold=last)

    # ---------------- public slide builders ----------------
    def cover(self, title, subtitle=None, meta_lines=None):
        """Light full-bleed title slide (not page-numbered)."""
        s = self._slide()
        # white background with a left accent spine — same red accent language as
        # the interior pages (no heavy top band, keeps cover & body coordinated)
        self._rect(s, 0, 0, Inches(0.32), self.H, "accent")
        self._rect(s, 0, 0, self.W, Inches(0.05), "accent")
        # title — auto-shrink to keep it to ≤2 lines, then lay the underline,
        # subtitle and meta BELOW the title's actual height so a long (wrapping)
        # title never collides with the red underline.
        title_x = Inches(1.1)
        title_w_in = 11.2
        tsize = 40.0
        units = self._text_units(title)
        for _ in range(20):                       # shrink until it fits in 2 lines
            per_line = max(1.0, title_w_in * 72.0 / tsize)
            nlines = max(1, math.ceil(units / per_line))
            if nlines <= 2 or tsize <= 26:
                break
            tsize = max(26.0, tsize * 0.94)
        per_line = max(1.0, title_w_in * 72.0 / tsize)
        nlines = max(1, math.ceil(units / per_line))
        title_top = 2.15
        title_h = nlines * tsize * 1.32 / 72.0
        tf = self._textbox(s, title_x, Inches(title_top), Inches(title_w_in),
                           Inches(title_h + 0.1), MSO_ANCHOR.TOP)
        r = tf.paragraphs[0].add_run(); r.text = title
        self._run(r, tsize, "primary", bold=True)
        # accent underline — placed just under the real title bottom
        y = title_top + title_h + 0.12
        self._rect(s, Inches(1.13), Inches(y), Inches(3.4), Inches(0.06), "accent")
        y += 0.18
        if subtitle:
            tf = self._textbox(s, Inches(1.13), Inches(y), Inches(11), Inches(0.8))
            r = tf.paragraphs[0].add_run(); r.text = subtitle
            self._run(r, 20, "grey")
            y += 1.0
        else:
            y += 0.55
        if meta_lines:
            tf = self._textbox(s, Inches(1.13), Inches(y), Inches(11), Inches(1.4))
            for i, ln in enumerate(meta_lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(6)
                r = p.add_run(); r.text = str(ln)
                self._run(r, 15, "grey")
        self._footer(s)
        return s

    def content_slide(self, title, cards, subtitle=None, takeaway=None):
        """Auto-grid of 1–6 cards. Each card: {title, bullets, accent}."""
        _warn_wrong_builder(title, "content_slide", "value_slide / design_slide")
        s = self._slide()
        self._header(s, title, subtitle, takeaway)
        n = len(cards)
        cols = 1 if n == 1 else (2 if n in (2, 4) else 3)
        rows = (n + cols - 1) // cols
        gap = Inches(0.3)
        total_w = self.W - 2 * self.MARGIN
        cw = (total_w - gap * (cols - 1)) / cols
        total_h = self.BODY_BOTTOM - self.BODY_TOP
        ch = (total_h - gap * (rows - 1)) / rows
        for i, c in enumerate(cards):
            rr, cc = divmod(i, cols)
            l = self.MARGIN + cc * (cw + gap)
            t = self.BODY_TOP + rr * (ch + gap)
            self._card(s, l, t, cw, ch,
                       c.get("title", ""), c.get("bullets", []),
                       # force one accent family — per-card colors are ignored on
                       # purpose so cards never become a multi-colored ("不纯粹") grid
                       "accent")
        return s

    def banded_slide(self, title, sections, subtitle=None, takeaway=None):
        """Vertical stack of full-width horizontal bars (横框), top to bottom.
        Each section is one bar: {title, bullets, accent}. Title sits in a left
        label column with a colored accent stripe; points fill the right column.
        Use this instead of content_slide when the brief wants sections shown
        as stacked horizontal frames rather than a card grid."""
        _warn_wrong_builder(title, "banded_slide", "value_slide / design_slide")
        s = self._slide()
        self._header(s, title, subtitle, takeaway)
        n = len(sections)
        gap = Inches(0.14)
        x = self.MARGIN
        w = self.W - 2 * self.MARGIN
        total_h = self.BODY_BOTTOM - self.BODY_TOP
        bh = (total_h - gap * (n - 1)) / n
        label_w = Inches(3.1)
        bsize = 13 if n <= 4 else 12
        for i, sec in enumerate(sections):
            t = self.BODY_TOP + i * (bh + gap)
            # force one accent family — per-section colors are ignored on purpose
            # so the horizontal bars never turn into multi-colored ("不纯粹") 横条
            accent = "accent"
            # band background + thin solid border + left accent stripe
            self._rect(s, x, t, w, bh, _tint(accent), line=_RULE, line_w=0.75)
            self._rect(s, x, t, Inches(0.12), bh, accent)
            # divider between label column and content column
            self._rect(s, x + label_w, t + Inches(0.10),
                       Inches(0.013), bh - Inches(0.20), _RULE)
            # left: section title
            tf = self._textbox(s, x + Inches(0.30), t,
                               label_w - Inches(0.45), bh, MSO_ANCHOR.MIDDLE)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = sec.get("title", "")
            self._run(r, 14, "dark", bold=True)
            # right: points
            tf = self._textbox(s, x + label_w + Inches(0.28), t,
                               w - label_w - Inches(0.55), bh, MSO_ANCHOR.MIDDLE)
            for j, ln in enumerate(sec.get("bullets", [])):
                pp = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                pp.space_after = Pt(2)
                r = pp.add_run(); r.text = "• " + str(ln)
                self._run(r, bsize, "dark")
        return s

    def bullets_slide(self, title, bullets, subtitle=None, takeaway=None):
        """Single text column. bullets: list of str or {text, level, accent, bold}."""
        s = self._slide()
        self._header(s, title, subtitle, takeaway)
        tf = self._textbox(s, self.MARGIN, self.BODY_TOP,
                           self.W - 2 * self.MARGIN,
                           self.BODY_BOTTOM - self.BODY_TOP, MSO_ANCHOR.MIDDLE)
        n = max(1, len(bullets))
        gap_pt = max(6, min(22, int(150 / n)))
        for i, b in enumerate(bullets):
            if isinstance(b, str):
                b = {"text": b}
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(gap_pt)
            p.space_before = Pt(gap_pt if i else 0)
            lvl = b.get("level", 0)
            mark = "■ " if lvl == 0 else ("• " if lvl == 1 else "– ")
            indent = "" if lvl == 0 else ("    " * lvl)
            r = p.add_run(); r.text = indent + mark + b["text"]
            self._run(r, 16 - min(lvl, 2) * 1.5, b.get("accent", "dark"),
                      bold=b.get("bold", lvl == 0))
        return s

    def table_slide(self, title, headers, rows, subtitle=None, takeaway=None,
                    col_widths=None, highlight_last=False):
        """headers: list[str]. rows: list[list[str]]. col_widths in inches.
        Row heights are sized to the REAL wrapped-text height of each cell
        (measured per column width), not a flat row/column-count average —
        so a row with a long cell gets the room it needs instead of being
        squeezed to the same height as a short row. If the measured content
        still doesn't fit the page at minimum font size, the table is left
        at its true (taller) height rather than force-compressed, so the
        Verification overflow check below can actually catch it."""
        s = self._slide()
        self._header(s, title, subtitle, takeaway)
        ncol = len(headers)
        nrow = len(rows) + 1
        total_w = self.W - 2 * self.MARGIN
        tbl = s.shapes.add_table(nrow, ncol, self.MARGIN, self.BODY_TOP,
                                 total_w, Inches(0.4)).table
        if col_widths:
            scale = float(total_w) / sum(Inches(w) for w in col_widths)
            col_w_emu = [int(Inches(w) * scale) for w in col_widths]
        else:
            col_w_emu = [int(total_w / ncol)] * ncol
        for ci, wv in enumerate(col_w_emu):
            tbl.columns[ci].width = Emu(wv)
        for ci, htext in enumerate(headers):
            tbl.cell(0, ci).text = str(htext)
        for ri, row in enumerate(rows, start=1):
            for ci, val in enumerate(row):
                tbl.cell(ri, ci).text = str(val)
        # Fonts follow the value/design page standard: header 15pt bold, body 13.5pt.
        # Dense tables shrink: many rows shrink the body; many columns (wide tables)
        # shrink both header and body so 11-column plans still fit horizontally.
        # The loop below refines this further using MEASURED cell content.
        body = 13.5 if nrow <= 9 else (12 if nrow <= 12 else 10.5)
        header = 15
        if ncol >= 9:
            body = min(body, 9.0); header = min(header, 11.0)
        elif ncol >= 7:
            body = min(body, 10.5); header = min(header, 12.5)

        # cell inner width in inches (minus ~0.1" left/right text-frame margins)
        col_w_in = [emu / 914400.0 - 0.2 for emu in col_w_emu]
        body_h_in = (self.BODY_BOTTOM - self.BODY_TOP) / 914400.0

        def row_heights_at(h_pt, b_pt):
            heights = [max(self._wrapped_rows(headers[ci], h_pt, col_w_in[ci])
                          for ci in range(ncol)) * h_pt * 1.35 / 72.0 + 0.12]
            for row in rows:
                heights.append(
                    max(self._wrapped_rows(row[ci], b_pt, col_w_in[ci])
                       for ci in range(ncol)) * b_pt * 1.30 / 72.0 + 0.10)
            return heights

        min_header, min_body = 8.0, 7.5
        h_pt, b_pt = float(header), float(body)
        heights = row_heights_at(h_pt, b_pt)
        for _ in range(14):
            if sum(heights) <= body_h_in or (h_pt <= min_header and b_pt <= min_body):
                break
            h_pt = max(min_header, h_pt * 0.93)
            b_pt = max(min_body, b_pt * 0.93)
            heights = row_heights_at(h_pt, b_pt)
        self._style_table(tbl, h_pt, b_pt, highlight_last=highlight_last)
        if sum(heights) > body_h_in:
            warnings.warn(
                "deckbuilder: table %r content does not fit the page even at "
                "minimum font size (%.1fpt/%.1fpt) — the table will overflow "
                "the canvas. Split the rows across multiple table_slide() "
                "calls or shorten cell text." % (title or "", h_pt, b_pt),
                stacklevel=2)
        else:
            # distribute leftover room proportionally so rows still fill the
            # page instead of leaving a gap at the bottom
            total_needed = sum(heights)
            slack = body_h_in - total_needed
            if slack > 0 and total_needed > 0:
                heights = [hh + slack * (hh / total_needed) for hh in heights]
        for ri, hh in enumerate(heights):
            tbl.rows[ri].height = Emu(int(hh * 914400))
        return s

    def flow_slide(self, title, stages, subtitle=None, takeaway=None, note=None,
                   lane_label=None):
        """Single horizontal pipeline with auto right-arrows.
        stages: list of {title, lines, change}. change=True → orange ★变更 badge."""
        s = self._slide()
        self._header(s, title, subtitle, takeaway)
        n = len(stages)
        gap = 0.40
        total_w = 13.333 - 2 * 0.4
        bw = (total_w - gap * (n - 1)) / n
        bh = 1.7
        y = 3.0 if note else 3.4
        if lane_label:
            tf = self._textbox(s, Inches(0.4), Inches(y - 0.45),
                               Inches(6), Inches(0.34))
            r = tf.paragraphs[0].add_run(); r.text = lane_label
            self._run(r, 13, "accent", bold=True)
        for i, st in enumerate(stages):
            x = 0.4 + i * (bw + gap)
            self._flowbox(s, Inches(x), Inches(y), Inches(bw), Inches(bh),
                          st.get("title", ""), st.get("lines", []),
                          change=st.get("change", False))
            if i < n - 1:
                ax = x + bw + 0.03
                self._arrow(s, Inches(ax), Inches(y + bh / 2 - 0.15),
                            Inches(gap - 0.06), Inches(0.30),
                            MSO_SHAPE.RIGHT_ARROW)
        if note:
            self._card(s, Inches(0.4), Inches(5.15), Inches(12.53), Inches(1.6),
                       "★ 说明 / Notes",
                       note if isinstance(note, list) else [note],
                       accent="grey")
        return s

    def layered_diagram_slide(self, title, layers, subtitle=None, takeaway=None,
                              connect=None, note=None):
        """Stacked horizontal layers (control plane / data plane style)."""
        s = self._slide()
        self._header(s, title, subtitle, takeaway)
        nlay = len(layers)
        top = 1.75
        bottom = 5.0 if note else 6.95          # leave room for the note card
        avail = bottom - top
        # All layers share ONE column grid sized to the widest layer, so boxes
        # in different planes line up vertically and cross-plane links stay clean.
        has_label = any(layer.get("label") for layer in layers)
        x0 = 1.75 if has_label else 0.4
        span = 13.0 - x0
        maxn = max(len(layer.get("nodes", [])) for layer in layers)
        gap = 0.4
        # content-measured box size
        all_nodes = [nd for layer in layers for nd in layer.get("nodes", [])]
        mw, mh, _, _ = self._measure_boxes(all_nodes)
        bw = min(mw, 2.8, (span - gap * (maxn - 1)) / maxn)
        lh = min(mh, 0.80)
        gw = bw * maxn + gap * (maxn - 1)
        gx0 = x0 + (span - gw) / 2                           # centred shared grid
        col_x = lambda c: gx0 + c * (bw + gap)
        band = avail / nlay                                 # even vertical bands
        pos = {}
        for li, layer in enumerate(layers):
            nodes = layer.get("nodes", [])
            nn = len(nodes)
            ly = top + li * band + (band - lh) / 2          # centre box in band
            label = layer.get("label")
            if label:
                tf = self._textbox(s, Inches(0.4), Inches(ly - 0.02),
                                   Inches(1.3), Inches(lh), MSO_ANCHOR.MIDDLE)
                r = tf.paragraphs[0].add_run(); r.text = label
                self._run(r, 11, "grey", bold=True)
            # column per node: explicit "col" if every node has one, else a
            # centred contiguous block within the shared grid.
            if nn and all("col" in nd for nd in nodes):
                cols = [nd["col"] for nd in nodes]
            else:
                start = (maxn - nn) // 2
                cols = [start + i for i in range(nn)]
            for ni, nd in enumerate(nodes):
                x = col_x(cols[ni])
                self._flowbox(s, Inches(x), Inches(ly), Inches(bw), Inches(lh),
                              nd.get("title", ""), nd.get("lines", []),
                              change=nd.get("change", False))
                pos[(li, ni)] = (x, ly, bw, lh)
                if ni > 0 and cols[ni] - cols[ni - 1] == 1:   # adjacent → arrow
                    ax = col_x(cols[ni - 1]) + bw + 0.03
                    self._arrow(s, Inches(ax), Inches(ly + lh / 2 - 0.13),
                                Inches(gap - 0.06), Inches(0.26),
                                MSO_SHAPE.RIGHT_ARROW)
        for a, b in (connect or []):
            (x1, y1, w1, h1) = pos[tuple(a)]
            (x2, y2, w2, h2) = pos[tuple(b)]
            down = y2 > y1
            # connect the two box centers (handles horizontal offset cleanly)
            p1 = (x1 + w1 / 2, (y1 + h1) if down else y1)
            p2 = (x2 + w2 / 2, y2 if down else (y2 + h2))
            self._connector(s, p1, p2, color="orange", direction="f", width=1.75)
        if note:
            self._card(s, Inches(0.4), Inches(5.15), Inches(12.53), Inches(1.6),
                       "★ 关键变更点 / Notes",
                       note if isinstance(note, list) else [note],
                       accent="grey")
        return s

    # ---------------- architecture diagram (labeled directional edges) ------
    def _connector(self, slide, p1, p2, color="grey", direction="f", width=1.5):
        from pptx.enum.shapes import MSO_CONNECTOR
        from pptx.oxml.ns import qn
        (x1, y1), (x2, y2) = p1, p2
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
            Inches(x2), Inches(y2))
        conn.line.color.rgb = _color(color)
        conn.line.width = Pt(width)
        conn.shadow.inherit = False
        ln = conn.line._get_or_add_ln()
        for tag, on in (("headEnd", direction in ("b", "both")),
                        ("tailEnd", direction in ("f", "both"))):
            if on:
                e = ln.makeelement(qn("a:" + tag), {})
                e.set("type", "triangle"); e.set("w", "med"); e.set("len", "med")
                ln.append(e)
        return conn

    def _edge_label(self, slide, mx, my, text, accent="grey"):
        w = max(0.6, min(2.0, self._text_width(text, 8.5) + 0.18)); h = 0.22
        self._rect(slide, Inches(mx - w / 2), Inches(my - h / 2),
                   Inches(w), Inches(h), "white", line=accent, line_w=0.5,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = self._textbox(slide, Inches(mx - w / 2), Inches(my - h / 2),
                           Inches(w), Inches(h), MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text; self._run(r, 8.5, accent, bold=True)

    @staticmethod
    def _anchor_points(a, b):
        ax, ay, aw, ah = a; bx, by, bw, bh = b
        acx, acy = ax + aw / 2, ay + ah / 2
        bcx, bcy = bx + bw / 2, by + bh / 2
        dx, dy = bcx - acx, bcy - acy
        if abs(dx) >= abs(dy):
            return ((ax + aw, acy), (bx, bcy)) if dx >= 0 else ((ax, acy), (bx + bw, bcy))
        return ((acx, ay + ah), (bcx, by)) if dy >= 0 else ((acx, ay), (bcx, by + bh))

    def architecture_slide(self, title, nodes, edges, subtitle=None, takeaway=None, note=None):
        """High-level architecture / module-interaction diagram."""
        s = self._slide()
        self._header(s, title, subtitle, takeaway)
        ncols = max(n["col"] for n in nodes) + 1
        nrows = max(n["row"] for n in nodes) + 1
        top = 1.75
        bottom = 5.35 if note else 6.85
        x0 = 0.4
        span = 13.333 - 2 * x0
        avail_h = bottom - top
        # content-measured box size
        cw, ch, _, _ = self._measure_boxes(nodes)
        max_cw = (span - 0.55 * (ncols - 1)) / ncols
        cw = min(cw, max_cw, 2.8)
        max_ch = (avail_h - 0.4 * (nrows - 1)) / nrows
        ch = min(ch, max_ch, 0.80)
        hgap = (span - cw * ncols) / (ncols - 1) if ncols > 1 else 0
        hgap = min(hgap, 2.0)
        grid_w = cw * ncols + hgap * (ncols - 1)
        xstart = x0 + (span - grid_w) / 2
        vgap = (avail_h - ch * nrows) / (nrows - 1) if nrows > 1 else 0
        vgap = min(vgap, 1.0)
        stack_h = ch * nrows + vgap * (nrows - 1)
        ystart = top + (avail_h - stack_h) / 2
        pos = {}
        for n in nodes:
            x = xstart + n["col"] * (cw + hgap)
            y = ystart + n["row"] * (ch + vgap)
            self._flowbox(s, Inches(x), Inches(y), Inches(cw), Inches(ch),
                          n.get("title", ""), n.get("lines", []),
                          change=n.get("change", False))
            pos[n["id"]] = (x, y, cw, ch)
        for e in edges:
            p1, p2 = self._anchor_points(pos[e["from"]], pos[e["to"]])
            col = e.get("accent", "grey")
            self._connector(s, p1, p2, color=col, direction=e.get("dir", "f"))
            if e.get("label"):
                self._edge_label(s, (p1[0] + p2[0]) / 2, (p1[1] + p2[1]) / 2,
                                  e["label"], accent=col)
        if note:
            self._card(s, Inches(0.4), Inches(5.55), Inches(12.53), Inches(1.25),
                       "★ 模块交互说明 / Notes",
                       note if isinstance(note, list) else [note],
                       accent="grey")
        return s

    # ---------------- two-column value / design slides ----------------
    def _plain_header(self, slide, title, title_pt=28, takeaway=None,
                      numbered=True):
        """Title block for the two-column pages — like `_header` but the title
        size is configurable (default 28pt bold YaHei) and `takeaway` is OPTIONAL
        (these pages lead with their section content, not a one-line verdict).
        Keeps the page number (top-right), accent underline, and logo."""
        x = self.MARGIN
        if numbered:
            self._page += 1
            tf = self._textbox(slide, self.W - Inches(1.4), Inches(0.34),
                               Inches(0.9), Inches(0.5), MSO_ANCHOR.MIDDLE)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
            r = p.add_run(); r.text = "%02d" % self._page
            self._run(r, 14, "grey", bold=True)
        title_w = self.W - x - Inches(1.5)
        tf = self._textbox(slide, x, Inches(0.36), title_w, Inches(0.74),
                           MSO_ANCHOR.MIDDLE)
        tf.word_wrap = False
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = title
        tsize = self._fit_size(title, title_pt, title_w / 914400.0, minsize=18)
        # title — Microsoft YaHei, bold, 28pt (auto-shrinks only if too long)
        self._run(r, tsize, "primary", bold=True)
        self._rect(slide, x, Inches(1.16), Inches(2.4), Inches(0.055), "accent")
        self._rect(slide, x + Inches(2.4), Inches(1.175),
                   self.W - x - Inches(2.4) - Inches(0.5), Inches(0.025), _RULE)
        if takeaway and str(takeaway).strip():
            self._rect(slide, x, Inches(1.335), Inches(0.09), Inches(0.09),
                       "accent")
            tf = self._textbox(slide, x + Inches(0.17), Inches(1.235),
                               self.W - x - Inches(0.77), Inches(0.34))
            tf.word_wrap = False
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = takeaway
            tw = self.W - x - Inches(0.77)
            ts = self._fit_size(takeaway, 14, tw / 914400.0, minsize=11)
            self._run(r, ts, "accent", bold=True)
        self._footer(slide)

    def _place_image(self, slide, l, t, w, h, image, placeholder):
        """Right-column visual: fit `image` into the (l,t,w,h) box preserving its
        aspect ratio and centering it; if no usable image is given, draw a light
        bordered placeholder so the user knows where to paste a 场景图 / 架构图."""
        if image and os.path.exists(image):
            pic = slide.shapes.add_picture(image, l, t)
            nw, nh = int(pic.width), int(pic.height)
            if nw > 0 and nh > 0:
                scale = min(float(w) / nw, float(h) / nh)
                pic.width = int(nw * scale)
                pic.height = int(nh * scale)
                pic.left = int(l + (int(w) - pic.width) / 2)
                pic.top = int(t + (int(h) - pic.height) / 2)
            return pic
        # placeholder card
        self._rect(slide, l, t, w, h, "white", line=_BOX_LINE, line_w=1.0,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = self._textbox(slide, l + Inches(0.2), t, w - Inches(0.4), h,
                           MSO_ANCHOR.MIDDLE)
        for i, ln in enumerate(placeholder if isinstance(placeholder, list)
                               else [placeholder]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(4)
            r = p.add_run(); r.text = str(ln)
            self._run(r, 13 if i == 0 else 11, "grey", bold=(i == 0))

    def _mini_layered_diagram(self, slide, l, t, w, h, layers, caption=None):
        """Compact vertical layered diagram sized to fit an arbitrary box — used
        to embed a framework 框图 in a slide's right column. Each layer is a band
        with a small grey label on the left and a row of real rounded boxes;
        adjacent layers are joined by a centered grey down-arrow. Coordinates are
        all derived from (l,t,w,h) so the whole stack is fully contained."""
        def I(v):
            return Emu(int(v))
        n = max(1, len(layers))
        cap_h = Inches(0.40) if caption else 0
        avail_h = h - cap_h
        g = Inches(0.20)                                   # inter-band gap
        band_h = (avail_h - g * (n - 1)) / n
        label_w = Inches(0.50)
        box_l = l + label_w
        box_area_w = w - label_w
        inner_gap = Inches(0.12)
        cx = box_l + box_area_w / 2                         # arrow centre-x
        for i, layer in enumerate(layers):
            by = t + i * (band_h + g)
            lbl = layer.get("label")
            if lbl:
                tf = self._textbox(slide, I(l), I(by), I(label_w - Inches(0.04)),
                                   I(band_h), MSO_ANCHOR.MIDDLE)
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = lbl
                self._run(r, 8.5, "grey", bold=True)
            nodes = layer.get("nodes", [])
            nn = max(1, len(nodes))
            bw = (box_area_w - inner_gap * (nn - 1)) / nn
            for j, nd in enumerate(nodes):
                bx = box_l + j * (bw + inner_gap)
                self._flowbox(slide, I(bx), I(by), I(bw), I(band_h),
                              nd.get("title", ""), nd.get("lines", []),
                              change=nd.get("change", False),
                              title_pt=10, sub_pt=8.5)
            if i < n - 1:                                   # down-arrow to next band
                ay = by + band_h + Inches(0.015)
                self._arrow(slide, I(cx - Inches(0.085)), I(ay),
                            I(Inches(0.17)), I(g - Inches(0.03)),
                            MSO_SHAPE.DOWN_ARROW, color="grey")
        if caption:
            tf = self._textbox(slide, I(l), I(t + avail_h + Inches(0.03)),
                               I(w), I(cap_h))
            for k, ln in enumerate(caption if isinstance(caption, list)
                                   else [caption]):
                p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
                p.space_after = Pt(1)
                r = p.add_run(); r.text = str(ln)
                self._run(r, 8.5, "grey")

    @staticmethod
    def _split_label(text):
        """Split a body line into (label, separator, rest) at the first Chinese
        or ASCII colon, so the leading 小标题 (e.g. '用户痛点') can be bolded.
        Only treats it as a label when the colon is reasonably early (≤14 units),
        otherwise returns ('', '', '') meaning 'no label'."""
        for sep in ("：", ":"):
            i = text.find(sep)
            if 0 < i <= 14:
                return text[:i], sep, text[i + 1:]
        return "", "", ""

    @classmethod
    def _fit_two_col(cls, sections, col_w_in, col_h_in,
                     head_pt=15.0, body_pt=13.5, min_head=11.0, min_body=9.5):
        """Shrink the left-column heading+body fonts together until the estimated
        wrapped height of all sections fits col_h_in. Returns (head_pt, body_pt)."""
        if col_w_in <= 0 or col_h_in <= 0:
            return head_pt, body_pt

        def total_height(hp, bp):
            gap_before = 12.0 * (bp / 13.5)
            h = 0.0
            for si, sec in enumerate(sections):
                if si > 0:
                    h += gap_before / 72.0
                # heading (rarely wraps, but count it)
                head = str(sec.get("heading", ""))
                hrows = max(1, math.ceil(cls._text_units(head) /
                                         max(1.0, col_w_in * 72.0 / hp)))
                h += hrows * hp * 1.3 / 72.0 + 5.0 / 72.0
                for ln in sec.get("lines", []):
                    units = cls._text_units("• " + str(ln))
                    rows = max(1, math.ceil(units /
                                            max(1.0, col_w_in * 72.0 / bp)))
                    h += rows * bp * 1.32 / 72.0 + 4.0 / 72.0
            return h

        hp, bp = head_pt, body_pt
        for _ in range(24):
            if total_height(hp, bp) <= col_h_in:
                break
            if hp <= min_head and bp <= min_body:
                break
            hp = max(min_head, hp * 0.96)
            bp = max(min_body, bp * 0.96)
        return hp, bp

    def _two_col_sections(self, title, sections, image=None,
                          image_placeholder=None, takeaway=None, title_pt=28,
                          diagram=None, diagram_caption=None):
        """Render a page with a left text column of stacked sections and a right
        column holding a visual (a framework 框图 via `diagram=`, else an image or
        a placeholder).
        sections: list of {heading, lines, head_color, body_color}.
        Section heading → 15pt bold; body lines → 13.5pt, colored per section."""
        s = self._slide()
        self._plain_header(s, title, title_pt=title_pt, takeaway=takeaway)
        top = self.BODY_TOP + (Inches(0.18) if takeaway else 0)
        bottom = self.BODY_BOTTOM
        left_x = self.MARGIN
        left_w = Inches(6.9)
        gap = Inches(0.4)
        right_x = left_x + left_w + gap
        right_w = self.W - self.MARGIN - right_x
        # faint divider between the two columns
        self._rect(s, left_x + left_w + gap / 2, top + Inches(0.05),
                   Inches(0.013), bottom - top - Inches(0.10), _RULE)
        # left column — one text frame, sections as paragraph groups.
        # Auto-shrink heading+body fonts TOGETHER so long content never spills
        # past the page bottom: estimate the wrapped height at the default sizes
        # and scale down until it fits the available column height.
        col_h_in = (bottom - top) / 914400.0
        col_w_in = left_w / 914400.0 - 0.1                # minus bullet indent
        head_pt, body_pt = self._fit_two_col(sections, col_w_in, col_h_in)
        gap_before = 12.0 * (body_pt / 13.5)
        tf = self._textbox(s, left_x, top, left_w, bottom - top)
        first = True
        for sec in sections:
            head = sec.get("heading", "")
            hc = sec.get("head_color", "dark")
            bc = sec.get("body_color", "dark")
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.space_before = Pt(0 if first else gap_before)
            p.space_after = Pt(5)
            r = p.add_run(); r.text = head
            self._run(r, head_pt, hc, bold=True)       # heading bold
            first = False
            for ln in sec.get("lines", []):
                pp = tf.add_paragraph(); pp.space_after = Pt(4)
                # bold the leading label up to the first ：/: so 小标题
                # (用户痛点 / 用户场景 …) stands out from its body text.
                label, sep, rest = self._split_label(str(ln))
                r = pp.add_run(); r.text = "• "
                self._run(r, body_pt, bc, bold=False)
                if sep:
                    r = pp.add_run(); r.text = label + sep
                    self._run(r, body_pt, bc, bold=True)
                    if rest:
                        r = pp.add_run(); r.text = rest
                        self._run(r, body_pt, bc, bold=False)
                else:
                    r = pp.add_run(); r.text = str(ln)
                    self._run(r, body_pt, bc, bold=False)
        # right column — framework 框图, scene/architecture image, or placeholder
        if diagram:
            self._mini_layered_diagram(s, right_x + Inches(0.12), top,
                                       right_w - Inches(0.12), bottom - top,
                                       diagram, caption=diagram_caption)
        else:
            self._place_image(s, right_x, top, right_w, bottom - top,
                              image, image_placeholder)
        return s

    def value_slide(self, title="需求价值描述", background=None, features=None,
                    scope=None, image=None, takeaway=None):
        """需求价值描述页 — 左侧三段（背景 / 特性及价值点 / 需求收益（影响）范围），
        右侧贴场景图展示该特性的价值。
        - background: 背景 正文（蓝色 13.5pt）
        - features:   特性及价值点 正文（蓝色 13.5pt）—— 用户痛点 / 需求功能点 / 范围 /
                      用户场景 / 所带来的价值；toD 明确开发者适用范围与达成结果；
                      可量化目标（活跃数/好评率/留存）；产品范围（1+8+N）
        - scope:      需求收益（影响）范围 正文（蓝色 13.5pt）—— 使用范围 / 地区 /
                      具体产品；重点关注特性通用性（硬件依赖除外）
        - image:      右侧场景图路径（可选）"""
        sections = [
            {"heading": "背景", "lines": background or [],
             "head_color": "dark", "body_color": "blue"},
            {"heading": "特性及价值点", "lines": features or [],
             "head_color": "dark", "body_color": "blue"},
            {"heading": "需求收益（影响）范围", "lines": scope or [],
             "head_color": "dark", "body_color": "blue"},
        ]
        return self._two_col_sections(
            title, sections, image=image, takeaway=takeaway,
            image_placeholder=["场景图 / 价值示意",
                               "在此粘贴展示该特性价值的场景图"])

    def design_slide(self, title="需求设计方案", design=None, changes=None,
                     extra=None, image=None, takeaway=None, diagram=None,
                     diagram_caption=None):
        """需求设计方案页 — 左侧设计方案与变更点及影响，右侧贴架构图补充展示。
        设计方案可有多种、可多页输出（每页各调用一次本方法）。
        - design:  设计方案设计重点（交互模块及如何达成需求的规格等）
        - changes: 变更点及影响（数据结构变更 / 外部接口变更 / 外部依赖分析 /
                   性能功耗评估 / 影响用户体验的关键 KPI 等）
        - extra:   追加段落 [{"heading","lines"}]，如「UI 示意」「裁剪已上线需求需求方意见」
        - image:   右侧架构图路径（可选）
        - diagram: 右侧内嵌框图 layers（与 layered_diagram_slide 同结构，优先于 image）
                   [{"label","nodes":[{"title","lines","change"}]}]
        - diagram_caption: 框图下方小字说明（可选）"""
        sections = [
            {"heading": "一、设计方案", "lines": design or [],
             "head_color": "dark", "body_color": "dark"},
            {"heading": "二、变更点及影响", "lines": changes or [],
             "head_color": "dark", "body_color": "dark"},
        ]
        for sec in (extra or []):
            sections.append({"heading": sec.get("heading", ""),
                             "lines": sec.get("lines", []),
                             "head_color": "dark", "body_color": "dark"})
        return self._two_col_sections(
            title, sections, image=image, takeaway=takeaway,
            diagram=diagram, diagram_caption=diagram_caption,
            image_placeholder=["架构图 / 补充示意",
                               "在此粘贴架构图补充展示"])

    # ------------------------------------------------------------------
    # ONE-STOP requirement-CHANGE-review deck (8 fixed pages).
    # A weaker model only fills a plain dict; this code owns page order,
    # which builder each page uses, the fixed 5-row 影响 / 8-row 风险 /
    # 11-col 交付 structure, and auto-fills 待评估 / TBD for missing fields.
    # This removes the #1 failure of "model picks the wrong builder".
    # ------------------------------------------------------------------
    _TBD = "待评估 / TBD"

    def _rr_lines(self, v):
        """Normalize a section value to a clean list[str], dropping empties.
        Falls back to [待评估 / TBD] when nothing usable remains, so a missing
        field renders as an explicit placeholder rather than a blank section."""
        if v is None:
            v = []
        elif isinstance(v, str):
            v = [v]
        lines = [str(x) for x in v if str(x).strip()]
        return lines if lines else [self._TBD]

    def _rr_cell(self, v):
        """A table cell string; blank/None → 待评估 / TBD (never fabricate)."""
        if v is None:
            return self._TBD
        s = str(v).strip()
        return s if s else self._TBD

    def requirement_review_deck(self, spec):
        """Build the full 8-page OpenHarmony 需求变更评审 deck from ONE dict.

        The page order, the builder used per page, and the fixed table
        structures (page 5 = 5 影响对象 rows, page 6 = 11-column plan, page 8 =
        8-row 风险 checklist) are all enforced HERE in code — the caller only
        supplies content and never chooses a builder or a layout. Every field is
        optional; anything missing renders as 待评估 / TBD rather than fabricated.

        spec = {
          # 1 封面
          "title": "<特性/需求名称> 需求变更评审",
          "subtitle": "<一句话副标题>",
          "meta_lines": ["<子系统/SIG>", "需求变更评审稿  2026-06-29"],
          # 2 需求价值描述 (value_slide, 左文右图)
          "value": {
              "background": [...], "features": [...], "scope": [...],
              "image": None,                 # 右侧场景图路径 (可选)
              "takeaway": None,              # 可选
          },
          # 3 需求设计方案 (design_slide, 左文右图 + 右侧框图). Pass a list to
          #   emit multiple 设计方案 pages (each dict = one page).
          "design": {
              "takeaway": "结论：…",
              "design": [...], "changes": [...], "extra": [ {...} ],
              "diagram": [ {"label","nodes":[{"title","lines","change"}]} ],
              "diagram_caption": [...], "image": None,
          },
          # 4 需求变更背景
          "background": {
              "takeaway": "结论：…",
              "requirement": "<原始需求/编号、特性概述与影响、使用场景>",
              "decision": "<原始被接纳时的 SIG 决策纪要>",
          },
          # 5 需求变更影响性分析 — 5 影响对象顺序固定, 只填每行分析文字
          "impact": {
              "takeaway": "结论：…",
              "north": "<北向应用开发者影响>",
              "south": "<南向开发者影响>",
              "distributed": "<分布式设备影响>",
              "system": "<系统开发者/跨子系统依赖影响>",
              "user": "<设备使用者 性能/功耗/功能/体验 影响>",
          },
          # 6 版本交付计划 — 11 列固定. items = 子需求列表, 一行一个; 自动加合计行
          "delivery": {
              "takeaway": "结论：…",
              "items": [
                {"domain","type","content","version","designer","loc",
                 "api","effort","pm","pipeline","reviewed"}, ...
              ],
          },
          # 7 兼容性分析
          "compatibility": {
              "takeaway": "结论：…",
              "involved": "<是否涉及应用兼容性>",
              "includes": "<机制/权限/API行为/其它 四类判定>",
              "plan": "<兼容性方案>",
              "adaptation": "<应用适配方案和计划>",
          },
          # 8 风险评估 — 8 行 checklist 固定, 只填每行结论
          "risk": {
              "takeaway": "结论：…",
              "perf": "...", "deps": "...", "security": "...", "legal": "...",
              "commitment": "...", "opensource": "...", "ai": "...",
              "privacy": "...",
          },
        }
        """
        spec = spec or {}
        C, T = self._rr_cell, self._TBD

        # 1 封面 --------------------------------------------------------
        self.cover(
            spec.get("title", "需求变更评审"),
            subtitle=spec.get("subtitle"),
            meta_lines=spec.get("meta_lines"),
        )

        # 2 需求价值描述 — value_slide (左文右图) -----------------------
        val = spec.get("value") or {}
        self.value_slide(
            title=val.get("title", "需求价值描述"),
            background=self._rr_lines(val.get("background")),
            features=self._rr_lines(val.get("features")),
            scope=self._rr_lines(val.get("scope")),
            image=val.get("image"),
            takeaway=val.get("takeaway"),
        )

        # 3 需求设计方案 — design_slide (可多页) ------------------------
        designs = spec.get("design") or {}
        if isinstance(designs, dict):
            designs = [designs]
        if not designs:
            designs = [{}]
        for dz in designs:
            self.design_slide(
                title=dz.get("title", "需求设计方案"),
                design=self._rr_lines(dz.get("design")),
                changes=self._rr_lines(dz.get("changes")),
                extra=dz.get("extra"),
                image=dz.get("image"),
                diagram=dz.get("diagram"),
                diagram_caption=dz.get("diagram_caption"),
                takeaway=dz.get("takeaway"),
            )

        # 4 需求变更背景 — table_slide ---------------------------------
        bg = spec.get("background") or {}
        self.table_slide(
            "四、需求变更背景",
            ["分项", "说明"],
            [
                ["需裁剪/变更的需求说明\n（重点说明 2C/2D 详细价值及产生影响）",
                 C(bg.get("requirement"))],
                ["原始被接纳时的决策纪要", C(bg.get("decision"))],
            ],
            takeaway=bg.get("takeaway") or "结论：" + T,
            col_widths=[2.6, 7])

        # 5 需求变更影响性分析 — 固定 5 行 -----------------------------
        im = spec.get("impact") or {}
        self.table_slide(
            "五、需求变更影响性分析",
            ["影响对象", "影响分析"],
            [
                ["北向应用开发者", C(im.get("north"))],
                ["南向开发者", C(im.get("south"))],
                ["分布式设备", C(im.get("distributed"))],
                ["系统开发者\n（跨子系统/部件依赖）", C(im.get("system"))],
                ["设备使用者\n（性能/功耗/功能/体验）", C(im.get("user"))],
            ],
            takeaway=im.get("takeaway") or "结论：" + T,
            col_widths=[2.4, 7])

        # 6 版本交付计划 — 固定 11 列, 子需求逐行 + 合计行 --------------
        dl = spec.get("delivery") or {}
        items = dl.get("items") or []
        headers = ["承接领域", "承接类型", "主要需求内容", "落地版本", "设计者",
                   "代码行数", "是否涉及API", "端到端工作量", "领域PM",
                   "管道是否满足", "工作量是否由领域PM审核OK"]
        keys = ["domain", "type", "content", "version", "designer", "loc",
                "api", "effort", "pm", "pipeline", "reviewed"]
        rows = []
        sum_loc = 0.0
        sum_eff = 0.0
        unknown_loc = 0
        unknown_eff = 0
        for it in items:
            it = it or {}
            rows.append([C(it.get(k)) for k in keys])
            n_loc = _num(it.get("loc"))
            if n_loc is None:
                unknown_loc += 1
            else:
                sum_loc += n_loc
            n_eff = _num(it.get("effort"))
            if n_eff is None:
                unknown_eff += 1
            else:
                sum_eff += n_eff
        if not rows:                               # never an empty plan table
            rows.append([C(None) for _ in keys])
            unknown_loc = unknown_eff = 1

        def _total_label(total, unknown, n_items):
            if n_items == 0 or unknown >= n_items:
                return T
            known = "%g" % total
            return known if unknown == 0 else "%s（另有 %d 项待评估）" % (known, unknown)

        total_loc = _total_label(sum_loc, unknown_loc, len(items) or 1)
        total_eff = _total_label(sum_eff, unknown_eff, len(items) or 1)
        rows.append(["合计（估算）", "—", "—", "—", "—", total_loc, "—",
                     total_eff, "—", "—", "—"])
        self.table_slide(
            "六、版本交付计划", headers, rows,
            takeaway=dl.get("takeaway") or "结论：" + T,
            col_widths=[1.4, 0.9, 2.6, 1.2, 0.9, 1.0, 1.3, 1.3, 0.9, 1.0, 1.4],
            highlight_last=True)

        # 7 兼容性分析 — table_slide -----------------------------------
        cp = spec.get("compatibility") or {}
        self.table_slide(
            "七、兼容性分析",
            ["分项", "内容"],
            [
                ["是否涉及应用兼容性", C(cp.get("involved"))],
                ["兼容性包括", C(cp.get("includes"))],
                ["兼容性方案", C(cp.get("plan"))],
                ["应用适配方案和计划", C(cp.get("adaptation"))],
            ],
            takeaway=cp.get("takeaway") or "结论：" + T,
            col_widths=[2.2, 7])

        # 8 风险评估 — 固定 2×8 checklist ------------------------------
        rk = spec.get("risk") or {}
        self.table_slide(
            "八、风险评估",
            ["评估项", "结论 / 说明"],
            [
                ["对性能、功耗、RAM/ROM 是否有影响", C(rk.get("perf"))],
                ["是否存在其他依赖关系？", C(rk.get("deps"))],
                ["是否有安全风险", C(rk.get("security"))],
                ["是否涉及合法 / 合规问题？", C(rk.get("legal"))],
                ["是否涉及外部承诺？", C(rk.get("commitment"))],
                ["是否开源", C(rk.get("opensource"))],
                ["是否涉及 AI", C(rk.get("ai"))],
                ["隐私风险特性识别", C(rk.get("privacy"))],
            ],
            takeaway=rk.get("takeaway") or "结论：" + T,
            col_widths=[3.0, 6])
        return self

    def save(self, path):
        self.prs.save(path)
        return path


# --------------------------------------------------------------------------
# demo / smoke test
# --------------------------------------------------------------------------
if __name__ == "__main__":
    d = Deck()
    d.cover("Deckbuilder Demo", "Every slide type in one file",
            ["deckbuilder.py", "smoke test"])
    d.content_slide("Content slide — auto card grid", [
        {"title": "Type", "bullets": ["New feature", "Not a bugfix"]},
        {"title": "Effort", "bullets": ["~6.5 person-months"]},
        {"title": "Scope", "bullets": ["Code + tests only"]},
    ], takeaway="结论：新增能力，约 6.5 人月，仅代码与测试")
    d.table_slide("Table slide", ["Task", "Owner", "Effort"],
                  [["API", "TBD", "1.0"], ["Tests", "TBD", "1.0"],
                   ["Total", "—", "2.0"]], highlight_last=True,
                  takeaway="结论：合计约 2.0 人月，人员待指派")
    d.flow_slide("Flow slide — pipeline with change points", [
        {"title": "HID device", "lines": ["USB / BT"]},
        {"title": "normalize", "lines": ["resolve binding"], "change": True},
        {"title": "windows mgr", "lines": ["hit test"], "change": True},
        {"title": "dispatch", "lines": ["UDS"]},
    ], takeaway="结论：变更集中在 normalize 与窗口命中两处",
       note="Amber ★变更 badges mark components that change.")
    out = d.save("deckbuilder_demo.pptx")
    print("saved:", out, "slides:", len(d.prs.slides._sldIdLst))
